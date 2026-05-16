#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx", "lxml"]
# ///
"""Convert Mermaid diagrams to .pptx with native editable shapes."""

import argparse
import glob
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path

from lxml import etree as lxml_et
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

EXCLUDE_DIRS = frozenset({'_reference', '_test-run', 'node_modules', '.build', 'Derived', '.git', 'build'})
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)  # 16:9
SVG = 'http://www.w3.org/2000/svg'
S = f'{{{SVG}}}'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# Light-theme palette
SLIDE_BG     = RGBColor(0xFF, 0xFF, 0xFF)
NODE_FILL    = RGBColor(0xF5, 0xF5, 0xF7)
NODE_LINE    = RGBColor(0x33, 0x33, 0x33)
CLUSTER_FILL = RGBColor(0xFA, 0xFA, 0xFC)
CLUSTER_LINE = RGBColor(0x99, 0x99, 0x99)
EDGE_COLOR   = RGBColor(0x33, 0x33, 0x33)
NODE_TEXT    = RGBColor(0x11, 0x11, 0x22)
LABEL_TEXT   = RGBColor(0x33, 0x33, 0x33)

# Sequence-diagram extras
ACTOR_FILL   = RGBColor(0xEC, 0xEC, 0xFF)
ACTOR_LINE   = RGBColor(0x93, 0x70, 0xDB)
LIFELINE     = RGBColor(0x99, 0x99, 0x99)
NOTE_FILL    = RGBColor(0xFF, 0xF5, 0xAD)
NOTE_LINE    = RGBColor(0xAA, 0xAA, 0x33)
LOOP_LINE    = RGBColor(0x93, 0x70, 0xDB)
LABEL_BOX    = RGBColor(0xEC, 0xEC, 0xFF)

FONT_NAME      = 'Produkt ExtraLight'
FONT_SIZE      = Pt(18)
SEQ_ACTOR_SIZE = Pt(11)
SEQ_MSG_SIZE   = Pt(9)
SEQ_NOTE_SIZE  = Pt(9)
SEQ_LABEL_SIZE = Pt(9)
EDGE_WIDTH     = Pt(2)


# ── helpers ───────────────────────────────────────────────────────────────────

def find_repo_root() -> Path:
    p = Path(__file__).resolve().parent
    while p != p.parent:
        if (p / '.git').exists():
            return p
        p = p.parent
    return Path.cwd()


def extract_md_diagrams(md_path: Path) -> list[tuple[str, str]]:
    content = md_path.read_text(encoding='utf-8')
    return [
        (f'{md_path.stem}-{i:02d}', m.group(1))
        for i, m in enumerate(re.finditer(r'```mermaid\n(.*?)\n```', content, re.DOTALL), 1)
    ]


def discover(paths: list[Path], repo_root: Path) -> list[tuple[str, str]]:
    diagrams: list[tuple[str, str]] = []
    if paths:
        for p in paths:
            if p.suffix == '.mmd':
                diagrams.append((p.stem, p.read_text(encoding='utf-8')))
            elif p.suffix == '.md':
                diagrams.extend(extract_md_diagrams(p))
        return diagrams
    for dirpath_str, dirnames, filenames in os.walk(repo_root):
        dirnames[:] = sorted(d for d in dirnames if d not in EXCLUDE_DIRS and not d.startswith('.'))
        for fname in sorted(filenames):
            fp = Path(dirpath_str) / fname
            if fp.suffix == '.mmd':
                diagrams.append((fp.stem, fp.read_text(encoding='utf-8')))
            elif fp.suffix == '.md':
                diagrams.extend(extract_md_diagrams(fp))
    return diagrams


def parse_translate(t: str) -> tuple[float, float]:
    m = re.search(r'translate\(([^,)]+)(?:,\s*([^)]+))?\)', t or '')
    if not m:
        return 0.0, 0.0
    return float(m.group(1)), float(m.group(2) or 0)


def get_text(g: lxml_et._Element) -> str:
    for fo in g.iter(f'{S}foreignObject'):
        text = ''.join(fo.itertext()).strip().replace('\xa0', ' ')
        if text:
            return text
    for t in g.iter(f'{S}text'):
        text = ''.join(t.itertext()).strip()
        if text:
            return text
    return ''


def get_parent_translate(elem) -> tuple[float, float]:
    """Sum translate() transforms from all ancestor elements."""
    tx, ty = 0.0, 0.0
    parent = elem.getparent()
    while parent is not None:
        t = parent.get('transform', '')
        if t:
            dx, dy = parse_translate(t)
            tx += dx
            ty += dy
        parent = parent.getparent()
    return tx, ty


def path_endpoints(d: str) -> tuple[float, float, float, float]:
    nums = [float(n) for n in re.findall(r'-?(?:\d+\.?\d*|\.\d+)', d)]
    return (nums[0], nums[1], nums[-2], nums[-1]) if len(nums) >= 4 else (0, 0, 0, 0)


def _ensure_ln(shape) -> lxml_et._Element | None:
    """Return the <a:ln> child of the shape's spPr, creating it if absent."""
    el = shape._element
    spPr = el.find(f'{{{P_NS}}}spPr')
    if spPr is None:
        return None
    ln = spPr.find(f'{{{A_NS}}}ln')
    if ln is None:
        ln = lxml_et.SubElement(spPr, f'{{{A_NS}}}ln')
    return ln


def add_arrowhead(connector, style: str = 'arrow') -> None:
    """Add an arrowhead at the tail end of a connector."""
    ln = _ensure_ln(connector)
    if ln is None:
        return
    for old in ln.findall(f'{{{A_NS}}}tailEnd'):
        ln.remove(old)
    tail = lxml_et.SubElement(ln, f'{{{A_NS}}}tailEnd')
    tail.set('type', style)
    tail.set('w', 'med')
    tail.set('len', 'med')


def set_dashed(shape) -> None:
    """Mark a shape's line as dashed (prstDash val='dash')."""
    ln = _ensure_ln(shape)
    if ln is None:
        return
    for old in ln.findall(f'{{{A_NS}}}prstDash'):
        ln.remove(old)
    pd = lxml_et.SubElement(ln, f'{{{A_NS}}}prstDash')
    pd.set('val', 'dash')


# ── coordinate mapping ────────────────────────────────────────────────────────

class CoordMapper:
    def __init__(self, vx0: float, vy0: float, vw: float, vh: float, margin: float = 0.05):
        margin_w = SLIDE_W * margin
        margin_h = SLIDE_H * margin
        usable_w = SLIDE_W - 2 * margin_w
        usable_h = SLIDE_H - 2 * margin_h
        self.scale = min(usable_w / vw, usable_h / vh)
        self.ox = margin_w + (usable_w - vw * self.scale) / 2 - vx0 * self.scale
        self.oy = margin_h + (usable_h - vh * self.scale) / 2 - vy0 * self.scale

    def pt(self, x: float, y: float) -> tuple[int, int]:
        return int(x * self.scale + self.ox), int(y * self.scale + self.oy)

    def dim(self, d: float) -> int:
        return int(d * self.scale)


# ── rendering ─────────────────────────────────────────────────────────────────

def render_svg(content: str, out: Path) -> bool:
    with tempfile.NamedTemporaryFile(suffix='.mmd', mode='w', delete=False, encoding='utf-8') as f:
        f.write(content)
        tmp = Path(f.name)
    try:
        r = subprocess.run(
            ['npx', '--yes', '-p', '@mermaid-js/mermaid-cli', 'mmdc',
             '-i', str(tmp), '-o', str(out)],
            capture_output=True, text=True,
        )
        if r.returncode != 0:
            print(f'  mmdc: {r.stderr.strip()[:200]}', file=sys.stderr)
            return False
        return out.exists()
    finally:
        tmp.unlink(missing_ok=True)


def _set_run(run, text: str, size, color: RGBColor) -> None:
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = False


def _put_text(shape, lines: list[str], size, color: RGBColor, align=PP_ALIGN.CENTER) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _set_run(p.add_run(), line, size, color)


# ── flowchart renderer ───────────────────────────────────────────────────────

def render_flowchart(root, slide, mapper: CoordMapper) -> None:
    ns = {'svg': SVG}

    label_by_x: dict[int, str] = {}
    for g in root.xpath('.//svg:g[contains(@class,"cluster-label")]', namespaces=ns):
        tx, _ = parse_translate(g.get('transform', ''))
        ptx, _ = get_parent_translate(g)
        text = get_text(g)
        if text:
            label_by_x[round(tx + ptx)] = text

    def cluster_label(rect_x: float, rect_w: float) -> str:
        cx = rect_x + rect_w / 2
        if not label_by_x:
            return ''
        key = min(label_by_x, key=lambda k: abs(k - cx))
        return label_by_x[key] if abs(key - cx) < rect_w else ''

    cluster_label_boxes: list[tuple[int, int, int, str]] = []

    for g in root.xpath('.//svg:g[contains(@class,"cluster") and not(contains(@class,"cluster-label"))]',
                        namespaces=ns):
        r = g.find(f'{S}rect')
        if r is None:
            continue
        ptx, pty = get_parent_translate(g)
        rx = float(r.get('x', 0)) + ptx
        ry = float(r.get('y', 0)) + pty
        rw = float(r.get('width', 100))
        rh = float(r.get('height', 50))
        px, py = mapper.pt(rx, ry)
        pw, ph = mapper.dim(rw), mapper.dim(rh)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px, py, pw, ph)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CLUSTER_FILL
        shape.line.color.rgb = CLUSTER_LINE
        shape.line.width = Pt(1.5)
        shape.text_frame.text = ''
        label = cluster_label(rx, rw)
        if label:
            cluster_label_boxes.append((px, py, pw, label))

    for path in root.xpath('.//svg:path[contains(@class,"flowchart-link")]', namespaces=ns):
        d = path.get('d', '')
        if not d:
            continue
        x1, y1, x2, y2 = path_endpoints(d)
        px1, py1 = mapper.pt(x1, y1)
        px2, py2 = mapper.pt(x2, y2)
        conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, px1, py1, px2, py2)
        conn.line.color.rgb = EDGE_COLOR
        conn.line.width = EDGE_WIDTH
        add_arrowhead(conn)

    for g in root.xpath('.//svg:g[contains(@class,"node") and contains(@class,"default")]',
                        namespaces=ns):
        cx, cy = parse_translate(g.get('transform', ''))
        ptx, pty = get_parent_translate(g)
        cx += ptx
        cy += pty
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        w, h = 100.0, 40.0
        for child in g:
            tag = child.tag.replace(f'{{{SVG}}}', '')
            if tag == 'rect':
                w = float(child.get('width', 100))
                h = float(child.get('height', 40))
                break
            elif tag == 'polygon':
                pts = [tuple(map(float, p.split(','))) for p in child.get('points', '').split() if ',' in p]
                if pts:
                    xs, ys = zip(*pts)
                    w, h = max(xs) - min(xs), max(ys) - min(ys)
                shape_type = MSO_AUTO_SHAPE_TYPE.DIAMOND
                break
            elif tag in ('circle', 'ellipse'):
                r = float(child.get('r', 20))
                w = float(child.get('rx', r)) * 2
                h = float(child.get('ry', r)) * 2
                shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
                break

        px, py = mapper.pt(cx - w / 2, cy - h / 2)
        pw, ph = mapper.dim(w), mapper.dim(h)
        shape = slide.shapes.add_shape(shape_type, px, py, pw, ph)
        shape.fill.solid()
        shape.fill.fore_color.rgb = NODE_FILL
        shape.line.color.rgb = NODE_LINE
        shape.line.width = Pt(1.5)
        _put_text(shape, [get_text(g)], FONT_SIZE, NODE_TEXT)

    label_h = int(Pt(20))
    for lpx, lpy, lpw, text in cluster_label_boxes:
        tb = slide.shapes.add_textbox(lpx, lpy, lpw, label_h)
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), text, Pt(13), LABEL_TEXT)


# ── sequence-diagram renderer ────────────────────────────────────────────────

def _line_coords(line) -> tuple[float, float, float, float]:
    return (float(line.get('x1', 0)), float(line.get('y1', 0)),
            float(line.get('x2', 0)), float(line.get('y2', 0)))


def _text_xy(t) -> tuple[float, float]:
    return float(t.get('x', 0)), float(t.get('y', 0))


def _text_content(t) -> str:
    return ''.join(t.itertext()).strip()


def render_sequence(root, slide, mapper: CoordMapper) -> None:
    ns = {'svg': SVG}

    # 1. Control-structure frames (alt / loop / opt) — perimeter rect, dividers, labels
    for g in root.xpath('.//svg:g[@data-et="control-structure"]', namespaces=ns):
        perim, dividers = [], []
        for line in g.findall(f'{S}line'):
            if 'loopLine' not in (line.get('class') or ''):
                continue
            (perim if 'dasharray' not in (line.get('style') or '') else dividers).append(line)
        if len(perim) >= 4:
            xs, ys = [], []
            for l in perim:
                x1, y1, x2, y2 = _line_coords(l)
                xs += [x1, x2]; ys += [y1, y2]
            x0, y0, x1_, y1_ = min(xs), min(ys), max(xs), max(ys)
            px, py = mapper.pt(x0, y0)
            pw, ph = mapper.dim(x1_ - x0), mapper.dim(y1_ - y0)
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px, py, pw, ph)
            box.fill.background()
            box.line.color.rgb = LOOP_LINE
            box.line.width = Pt(1)
            box.text_frame.text = ''

        for d in dividers:
            x1, y1, x2, y2 = _line_coords(d)
            px1, py1 = mapper.pt(x1, y1)
            px2, py2 = mapper.pt(x2, y2)
            ln = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, px1, py1, px2, py2)
            ln.line.color.rgb = LOOP_LINE
            ln.line.width = Pt(0.75)
            set_dashed(ln)

        for poly in g.findall(f'{S}polygon'):
            if 'labelBox' not in (poly.get('class') or ''):
                continue
            pts = [tuple(map(float, p.split(','))) for p in (poly.get('points') or '').split() if ',' in p]
            if not pts:
                continue
            xs, ys = zip(*pts)
            px, py = mapper.pt(min(xs), min(ys))
            pw, ph = mapper.dim(max(xs) - min(xs)), mapper.dim(max(ys) - min(ys))
            lb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px, py, pw, ph)
            lb.fill.solid()
            lb.fill.fore_color.rgb = LABEL_BOX
            lb.line.color.rgb = LOOP_LINE
            lb.line.width = Pt(0.75)
            lb.text_frame.text = ''

        for t in g.findall(f'{S}text'):
            cls = t.get('class', '')
            if 'labelText' not in cls and 'loopText' not in cls:
                continue
            text = _text_content(t)
            if not text:
                continue
            tx, ty = _text_xy(t)
            est_w = max(60.0, 8.0 * len(text))
            est_h = 16.0
            px, py = mapper.pt(tx - est_w / 2, ty - est_h / 2)
            pw, ph = mapper.dim(est_w), mapper.dim(est_h)
            tb = slide.shapes.add_textbox(px, py, pw, ph)
            _put_text(tb, [text], SEQ_LABEL_SIZE, LABEL_TEXT)

    # 2. Notes — rect + grouped noteText lines
    for g in root.xpath('.//svg:g[@data-et="note"]', namespaces=ns):
        rect = g.find(f'{S}rect')
        if rect is None:
            continue
        rx = float(rect.get('x', 0))
        ry = float(rect.get('y', 0))
        rw = float(rect.get('width', 100))
        rh = float(rect.get('height', 40))
        px, py = mapper.pt(rx, ry)
        pw, ph = mapper.dim(rw), mapper.dim(rh)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px, py, pw, ph)
        shape.fill.solid()
        shape.fill.fore_color.rgb = NOTE_FILL
        shape.line.color.rgb = NOTE_LINE
        shape.line.width = Pt(0.75)
        lines = [_text_content(t) for t in g.findall(f'{S}text') if _text_content(t)]
        _put_text(shape, lines or [''], SEQ_NOTE_SIZE, NODE_TEXT)

    # 3. Lifelines — thin vertical gray lines between actor top + bottom boxes
    for line in root.xpath('.//svg:line[contains(@class,"actor-line")]', namespaces=ns):
        x1, y1, x2, y2 = _line_coords(line)
        px1, py1 = mapper.pt(x1, y1)
        px2, py2 = mapper.pt(x2, y2)
        conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, px1, py1, px2, py2)
        conn.line.color.rgb = LIFELINE
        conn.line.width = Pt(0.5)

    # 4. Actor boxes (both top and bottom) — pair each rect with its sibling text
    for rect in root.xpath('.//svg:rect[contains(@class,"actor")]', namespaces=ns):
        cls = rect.get('class', '')
        if 'actor-top' not in cls and 'actor-bottom' not in cls:
            continue
        rx = float(rect.get('x', 0))
        ry = float(rect.get('y', 0))
        rw = float(rect.get('width', 100))
        rh = float(rect.get('height', 40))
        px, py = mapper.pt(rx, ry)
        pw, ph = mapper.dim(rw), mapper.dim(rh)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px, py, pw, ph)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACTOR_FILL
        shape.line.color.rgb = ACTOR_LINE
        shape.line.width = Pt(1)

        parent = rect.getparent()
        label = ''
        if parent is not None:
            for t in parent.findall(f'{S}text'):
                if 'actor-box' in (t.get('class') or ''):
                    label = _text_content(t)
                    break
        _put_text(shape, [label], SEQ_ACTOR_SIZE, NODE_TEXT)

    # 5. Messages — connectors with arrowheads (solid messageLine0 / dashed messageLine1)
    for line in root.xpath('.//svg:line[contains(@class,"messageLine")]', namespaces=ns):
        cls = line.get('class', '')
        x1, y1, x2, y2 = _line_coords(line)
        px1, py1 = mapper.pt(x1, y1)
        px2, py2 = mapper.pt(x2, y2)
        conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, px1, py1, px2, py2)
        conn.line.color.rgb = EDGE_COLOR
        conn.line.width = Pt(1.25)
        if 'messageLine1' in cls:
            set_dashed(conn)
        add_arrowhead(conn)

    # 6. Message text labels
    for t in root.xpath('.//svg:text[contains(@class,"messageText")]', namespaces=ns):
        text = _text_content(t)
        if not text:
            continue
        tx, ty = _text_xy(t)
        est_w = max(80.0, 7.5 * len(text))
        est_h = 18.0
        px, py = mapper.pt(tx - est_w / 2, ty - est_h / 2)
        pw, ph = mapper.dim(est_w), mapper.dim(est_h)
        tb = slide.shapes.add_textbox(px, py, pw, ph)
        _put_text(tb, [text], SEQ_MSG_SIZE, LABEL_TEXT)


# ── SVG → pptx ────────────────────────────────────────────────────────────────

def diagram_kind(root) -> str:
    role = (root.get('aria-roledescription') or '').lower()
    if 'sequence' in role:
        return 'sequence'
    return 'flowchart'


def svg_to_pptx(svg_path: Path, out_pptx: Path) -> None:
    tree = lxml_et.parse(svg_path)
    root = tree.getroot()

    vb = (root.get('viewBox') or '').split()
    vx0, vy0, vw, vh = (float(v) for v in vb) if len(vb) == 4 else (0, 0, 800, 600)
    mapper = CoordMapper(vx0, vy0, vw, vh)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = next((l for l in prs.slide_layouts if l.name == 'Blank'), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLIDE_BG

    kind = diagram_kind(root)
    if kind == 'sequence':
        render_sequence(root, slide, mapper)
    else:
        render_flowchart(root, slide, mapper)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


# ── main ──────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument('sources', nargs='*', help='Files or globs (default: whole repo)')
    ap.add_argument('--list', action='store_true', help='Dry-run: list diagrams only')
    args = ap.parse_args()

    repo_root = find_repo_root()
    out_dir = repo_root / 'build' / 'pptx'

    paths: list[Path] = []
    for s in args.sources:
        expanded = glob.glob(s, recursive=True)
        paths.extend(Path(p) for p in (expanded or [s]))

    diagrams = discover(paths, repo_root)
    if not diagrams:
        print('No Mermaid diagrams found.')
        return

    if args.list:
        print(f'Found {len(diagrams)} diagram(s):')
        for stem, _ in diagrams:
            print(f'  → {out_dir / stem}.pptx')
        return

    ok = 0
    with tempfile.TemporaryDirectory() as tmp_dir:
        for stem, content in diagrams:
            out_pptx = out_dir / f'{stem}.pptx'
            tmp_svg = Path(tmp_dir) / f'{stem}.svg'
            print(f'  {stem}', end=' ... ', flush=True)
            if render_svg(content, tmp_svg):
                try:
                    svg_to_pptx(tmp_svg, out_pptx)
                    print('✓')
                    ok += 1
                except Exception as e:
                    print(f'✗  ({e})')
            else:
                print('✗  (render failed)')

    print(f'\n{ok}/{len(diagrams)} files written to {out_dir}')


if __name__ == '__main__':
    main()
